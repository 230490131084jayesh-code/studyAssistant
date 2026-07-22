"""Extracts raw text from uploaded documents and splits it into overlapping chunks."""
from pathlib import Path
from typing import List, Tuple

from pypdf import PdfReader
from pptx import Presentation
import docx

from app.config import CHUNK_SIZE, CHUNK_OVERLAP


def extract_text(file_path: Path) -> List[Tuple[str, str]]:
    """Returns a list of (page_label, text) tuples so we can cite locations."""
    ext = file_path.suffix.lower()

    if ext == ".pdf":
        return _extract_pdf(file_path)
    if ext == ".pptx":
        return _extract_pptx(file_path)
    if ext == ".docx":
        return _extract_docx(file_path)
    if ext in (".txt", ".md"):
        return _extract_text_file(file_path)

    raise ValueError(f"Unsupported file type: {ext}")


def _extract_pdf(file_path: Path) -> List[Tuple[str, str]]:
    reader = PdfReader(str(file_path))
    pages = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            pages.append((f"page {i}", text))
    return pages


def _extract_pptx(file_path: Path) -> List[Tuple[str, str]]:
    prs = Presentation(str(file_path))
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        parts.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text for c in row.cells]
                    parts.append(" | ".join(cells))
        text = "\n".join(parts)
        if text.strip():
            slides.append((f"slide {i}", text))
    return slides


def _extract_docx(file_path: Path) -> List[Tuple[str, str]]:
    doc = docx.Document(str(file_path))
    # Group paragraphs into pseudo-pages of ~40 paragraphs so citations stay useful
    paras = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            paras.append(" | ".join(c.text for c in row.cells))

    chunks = []
    group_size = 40
    for i in range(0, len(paras), group_size):
        group = paras[i:i + group_size]
        label = f"section {i // group_size + 1}"
        chunks.append((label, "\n".join(group)))
    return chunks


def _extract_text_file(file_path: Path) -> List[Tuple[str, str]]:
    text = file_path.read_text(errors="ignore")
    return [("document", text)]


def chunk_text(location_texts: List[Tuple[str, str]]) -> List[dict]:
    """Splits (location, text) pairs into overlapping chunks with metadata."""
    chunks = []
    for location, text in location_texts:
        text = text.strip()
        if not text:
            continue
        start = 0
        while start < len(text):
            end = start + CHUNK_SIZE
            piece = text[start:end].strip()
            if piece:
                chunks.append({"location": location, "text": piece})
            if end >= len(text):
                break
            start = end - CHUNK_OVERLAP
    return chunks
